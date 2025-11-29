import io
import csv
from typing import Optional
from fastapi import UploadFile

# Import optional dependencies for file parsing
try:
    import docx
except ImportError:
    docx = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pytesseract
    from pdf2image import convert_from_bytes
    from PIL import Image
except ImportError:
    pytesseract = None
    convert_from_bytes = None
    Image = None


import base64
from typing import Dict, Any, List, Union

# ... imports remain the same ...

async def process_file(file: UploadFile) -> Dict[str, Any]:
    """
    Process an uploaded file and return structured content for Multimodal LLMs.
    
    Returns:
        Dict with keys:
        - type: "text", "image", or "mixed"
        - content: The content (string for text, base64 for image, list for mixed)
        - media_type: MIME type (for images)
    """
    filename = file.filename.lower()
    content = await file.read()
    
    result = {
        "type": "text",
        "content": "",
        "filename": file.filename
    }
    
    if filename.endswith('.txt') or filename.endswith('.md') or filename.endswith('.py') or filename.endswith('.xml') or filename.endswith('.json') or filename.endswith('.js') or filename.endswith('.jsx') or filename.endswith('.ts') or filename.endswith('.tsx') or filename.endswith('.html') or filename.endswith('.css'):
        # Text files
        try:
            result["content"] = content.decode('utf-8')
        except UnicodeDecodeError:
            result["content"] = content.decode('utf-8', errors='replace')
            
    elif filename.endswith('.csv'):
        # CSV files
        try:
            decoded = content.decode('utf-8')
            f = io.StringIO(decoded)
            reader = csv.reader(f)
            rows = list(reader)
            result["content"] = "\n".join([",".join(row) for row in rows])
        except Exception as e:
            result["content"] = f"Error reading CSV: {str(e)}"
            
    elif filename.endswith('.pdf'):
        # PDF files - Hybrid approach
        # 1. Try to extract text first
        # 2. If text is sparse, convert pages to images for Vision
        if pypdf and convert_from_bytes:
            try:
                pdf_file = io.BytesIO(content)
                reader = pypdf.PdfReader(pdf_file)
                pages_text = []
                for page in reader.pages:
                    pages_text.append(page.extract_text())
                text_content = "\n\n".join(pages_text)
                
                # If text is sufficient, treat as text file
                if len(text_content.strip()) > 50:
                    result["content"] = text_content
                else:
                    # Low text content -> Treat as images (Vision)
                    print(f"DEBUG: PDF has low text content ({len(text_content.strip())} chars). Converting to images for Vision.")
                    images = convert_from_bytes(content)
                    
                    # Limit to first 5 pages to avoid huge payloads for now
                    max_pages = 5
                    mixed_content = []
                    
                    for i, image in enumerate(images[:max_pages]):
                        # Convert PIL image to base64
                        buffered = io.BytesIO()
                        image.save(buffered, format="JPEG")
                        img_str = base64.b64encode(buffered.getvalue()).decode()
                        
                        mixed_content.append({
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{img_str}"
                            }
                        })
                    
                    if len(images) > max_pages:
                        mixed_content.append({
                            "type": "text",
                            "text": f"\n[Note: PDF truncated. Showing first {max_pages} of {len(images)} pages.]"
                        })
                        
                    result["type"] = "mixed"
                    result["content"] = mixed_content
                    
            except Exception as e:
                result["content"] = f"Error reading PDF: {str(e)}"
        else:
             result["content"] = "PDF support not available (pypdf/pdf2image not installed)"

    elif filename.endswith('.docx'):
        # Word files
        if docx:
            try:
                doc_file = io.BytesIO(content)
                doc = docx.Document(doc_file)
                result["content"] = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                result["content"] = f"Error reading DOCX: {str(e)}"
        else:
            result["content"] = "DOCX support not available"
            
    elif filename.endswith('.pptx'):
        # PowerPoint files
        if Presentation:
            try:
                pptx_file = io.BytesIO(content)
                prs = Presentation(pptx_file)
                slides_text = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    slides_text.append("\n".join(slide_text))
                result["content"] = "\n\n--- Slide ---\n\n".join(slides_text)
            except Exception as e:
                result["content"] = f"Error reading PPTX: {str(e)}"
        else:
            result["content"] = "PPTX support not available"
            
    elif filename.endswith('.png') or filename.endswith('.jpg') or filename.endswith('.jpeg'):
        # Image files - Direct Vision support
        try:
            # Determine mime type
            mime_type = "image/png" if filename.endswith('.png') else "image/jpeg"
            
            # Encode to base64
            img_str = base64.b64encode(content).decode('utf-8')
            
            result["type"] = "image"
            result["content"] = img_str
            result["media_type"] = mime_type
            
        except Exception as e:
            result["type"] = "text"
            result["content"] = f"Error processing image: {str(e)}"
        
    else:
        result["content"] = f"[Unsupported file type: {file.filename}]"

    return result
